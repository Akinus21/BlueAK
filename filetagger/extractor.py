"""Extract text content from various file types for AI tagging."""
import subprocess
import logging
from pathlib import Path

logger = logging.getLogger("filetagger.extractor")


def get_category(ext: str, supported: dict) -> str:
    ext = ext.lower()
    for category, exts in supported.items():
        if ext in exts:
            return category
    return "other"


def extract_text(path: str, category: str, config: dict) -> str:
    """Extract text content from a file. Returns empty string on failure."""
    p = Path(path)
    max_chars = config.get("max_content_chars", 4000)

    try:
        if category == "documents":
            return _extract_document(p, max_chars)
        elif category == "images":
            if config.get("ocr_enabled", True):
                return _extract_image_ocr(p, max_chars)
            return ""
        elif category == "code":
            return _extract_text_file(p, max_chars)
        elif category == "audio":
            if config.get("whisper_enabled", False):
                return _extract_audio_whisper(p, max_chars)
            return f"Audio file: {p.name}"
        elif category == "video":
            return f"Video file: {p.name}"
        else:
            # Try plain text read as fallback
            return _extract_text_file(p, max_chars)
    except Exception as e:
        logger.warning(f"Extraction failed for {path}: {e}")
        return ""


def _extract_document(p: Path, max_chars: int) -> str:
    ext = p.suffix.lower()

    if ext == ".pdf":
        return _pdf_text(p, max_chars)
    elif ext in (".doc", ".docx"):
        return _docx_text(p, max_chars)
    elif ext in (".odt",):
        return _odt_text(p, max_chars)
    elif ext in (".xlsx", ".xls"):
        return _xlsx_text(p, max_chars)
    elif ext in (".pptx", ".ppt"):
        return _pptx_text(p, max_chars)
    elif ext in (".txt", ".md", ".rtf", ".csv"):
        return _extract_text_file(p, max_chars)
    return ""


def _pdf_text(p: Path, max_chars: int) -> str:
    try:
        import pdfminer.high_level
        text = pdfminer.high_level.extract_text(str(p))
        return (text or "")[:max_chars]
    except ImportError:
        pass
    # Fallback: pdftotext CLI
    try:
        result = subprocess.run(["pdftotext", str(p), "-"],
                                capture_output=True, text=True, timeout=30)
        return result.stdout[:max_chars]
    except Exception:
        return ""


def _docx_text(p: Path, max_chars: int) -> str:
    try:
        from docx import Document
        doc = Document(str(p))
        text = "\n".join(para.text for para in doc.paragraphs)
        return text[:max_chars]
    except Exception:
        return ""


def _odt_text(p: Path, max_chars: int) -> str:
    try:
        result = subprocess.run(["odt2txt", str(p)],
                                capture_output=True, text=True, timeout=30)
        return result.stdout[:max_chars]
    except Exception:
        return ""


def _xlsx_text(p: Path, max_chars: int) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                line = " ".join(str(c) for c in row if c is not None)
                if line.strip():
                    texts.append(line)
        return "\n".join(texts)[:max_chars]
    except Exception:
        return ""


def _pptx_text(p: Path, max_chars: int) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)[:max_chars]
    except Exception:
        return ""


def _extract_text_file(p: Path, max_chars: int) -> str:
    try:
        text = p.read_text(encoding="utf-8", errors="replace")
        return text[:max_chars]
    except Exception:
        return ""


def _extract_image_ocr(p: Path, max_chars: int) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(str(p))
        text = pytesseract.image_to_string(img)
        return text[:max_chars]
    except ImportError:
        logger.debug("pytesseract/PIL not installed, skipping OCR")
        return ""
    except Exception as e:
        logger.debug(f"OCR failed for {p}: {e}")
        return ""


def _extract_audio_whisper(p: Path, max_chars: int) -> str:
    try:
        import whisper
        model = whisper.load_model("base")
        result = model.transcribe(str(p))
        return result["text"][:max_chars]
    except ImportError:
        logger.debug("whisper not installed, skipping audio transcription")
        return ""
    except Exception as e:
        logger.debug(f"Whisper failed for {p}: {e}")
        return ""
